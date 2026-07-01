from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def apply_dark_theme(slide):
    # Set background to a dark purple/blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900

def style_title(title_shape):
    if title_shape:
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.name = 'Arial'

def style_body(body_shape):
    if body_shape:
        for p in body_shape.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(226, 232, 240) # Slate 200
            p.font.size = Pt(18)
            p.font.name = 'Arial'

def add_slide(prs, layout_idx, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    apply_dark_theme(slide)
    
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        style_title(slide.shapes.title)
        
    if len(slide.placeholders) > 1 and body_text:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        if isinstance(body_text, list):
            for i, point in enumerate(body_text):
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                if point.startswith("[MANDATORY VISUAL]"):
                    p.font.color.rgb = RGBColor(147, 197, 253) # Blue 300
                    p.font.bold = True
        else:
            p = tf.add_paragraph()
            p.text = body_text
            
        style_body(body_shape)
    
    return slide

prs = Presentation()

# Slide 1: Title
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
apply_dark_theme(title_slide)
title_shape = title_slide.shapes.title
subtitle_shape = title_slide.placeholders[1]

title_shape.text = "INDIA.RUNS - Track 1"
subtitle_shape.text = "AI Systems & Workflow Innovation Challenge\nTeam CortexX"
style_title(title_shape)
style_body(subtitle_shape)
subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Slide 2: Team Info
add_slide(prs, 1, "Team & Problem Statement", [
    "Team Name: CortexX Engineers",
    "Team Members: Antigravity AI Assistant & Developer",
    "Problem Statement: Bridging the gap between passive AI chat and active, autonomous software engineering within the Redrob ecosystem."
])

# Slide 3: Problem Definition
add_slide(prs, 1, "Problem Definition", [
    "What problem are you solving?",
    "  - Developer friction in building new workflows and analyzing complex system logs manually.",
    "Who experiences this problem?",
    "  - Engineers and developers building the Redrob ecosystem.",
    "Why is the current approach insufficient?",
    "  - Current AI tools are passive (AI-assisted). They require human intervention to execute commands, read files, and synthesize logs."
])

# Slide 4: Opportunity & Vision
add_slide(prs, 1, "Opportunity & Vision", [
    "Why is this an important opportunity?",
    "  - Automating the developer loop massively accelerates time-to-market for new Redrob features and internal workflows.",
    "What future state are you enabling?",
    "  - A state where developers act as managers of autonomous AI agents (CortexX) that write, debug, and deploy code independently."
])

# Slide 5: Solution Overview
add_slide(prs, 1, "Solution Overview", [
    "What is your proposed solution?",
    "  - CortexX: A Multi-Agent Autonomous Engineering System with a Glassmorphism UI and FastAPI backend.",
    "What makes it AI-native rather than AI-assisted?",
    "  - CortexX has physical agency. It uses LangChain ReAct loops to read files, edit code inline, and execute terminal commands autonomously.",
    "Which existing Redrob capability does it build upon?",
    "  - Builds upon Redrob's productivity and workflow ecosystem, serving as the ultimate meta-tool for engineering."
])

# Slide 6: User Journey / Workflow Diagram
add_slide(prs, 1, "User Journey / Workflow Diagram", [
    "How does a user interact with the solution?",
    "  - Through a Next.js chat interface or 'Project Builder' mode.",
    "How does information flow?",
    "  - User -> Next.js -> FastAPI -> LangChain Router -> Specialized Agents -> Local OS -> User.",
    "Where does it integrate?",
    "  - Integrates directly into Redrob's developer tooling suite.",
    "",
    "[MANDATORY VISUAL]: (Placeholder for Workflow Diagram)"
])

# Slide 7: AI Logic & Decision Flow
add_slide(prs, 1, "AI Logic & Decision Flow", [
    "Where does AI intervene?",
    "  - At every step—from intent parsing to code generation and terminal execution.",
    "How are decisions made?",
    "  - Using Google's Gemini models in a ReAct loop (Reasoning + Acting).",
    "How do agents interact?",
    "  - The MultiModelRouter delegates tasks to specialized sub-agents (CodeIntel, Debug, DevOps).",
    "",
    "[MANDATORY VISUAL]: (Placeholder for AI Flow Diagram)"
])

# Slide 8: System Architecture
add_slide(prs, 1, "System Architecture", [
    "What components make up the system?",
    "  - Next.js Frontend, FastAPI Backend, LangChain, Google Gemini API, Local File System Tools.",
    "How do services interact?",
    "  - REST API and SSE (Server-Sent Events) streaming for real-time agent thought processing.",
    "Which Redrob systems are leveraged?",
    "  - Redrob's authentication, productivity tracking, and deployment pipelines.",
    "",
    "[MANDATORY VISUAL]: (Placeholder for Architecture Diagram)"
])

# Slide 9: Data, Context & Intelligence Layer
add_slide(prs, 1, "Data, Context & Intelligence Layer", [
    "What data powers the solution?",
    "  - Local codebase files, terminal outputs, and live CI/CD logs.",
    "How is context retrieved, stored, or utilized?",
    "  - Semantic search (Knowledge Agent) and direct file reading tools.",
    "How does existing Redrob context improve the experience?",
    "  - By embedding Redrob's specific engineering guidelines directly into the agent's system prompt.",
    "",
    "[MANDATORY VISUAL]: (Placeholder for Data Flow Diagram)"
])

# Slide 10: Scalability & Technical Feasibility
add_slide(prs, 1, "Scalability & Technical Feasibility", [
    "How would this be implemented?",
    "  - Dockerized containers for isolated, safe terminal execution (sandbox).",
    "How does the system scale?",
    "  - Stateless FastAPI backend allows horizontal scaling; LLMs handle heavy lifting via cloud APIs.",
    "What technical challenges exist?",
    "  - Context window limits on massive codebases and preventing destructive terminal commands."
])

# Slide 11: Redrob Ecosystem Integration
add_slide(prs, 1, "Redrob Ecosystem Integration", [
    "Which existing Redrob capabilities are being leveraged?",
    "  - Redrob internal workflows and productivity tools.",
    "What new capability does your solution introduce?",
    "  - Autonomous Project Building and Self-Healing Code.",
    "How does it strengthen the overall ecosystem?",
    "  - Dramatically reduces internal engineering overhead.",
    "What additional opportunities become possible?",
    "  - Exposing this to end-users as a no-code AI workflow builder.",
    "",
    "[MANDATORY VISUAL]: (Placeholder for Ecosystem Integration Diagram)"
])

# Slide 12: Impact & Success Metrics
add_slide(prs, 1, "Impact & Success Metrics", [
    "What measurable outcomes are expected?",
    "  - 40% reduction in time-to-resolve for complex bugs; 3x faster scaffolding of new microservices.",
    "How will success be tracked?",
    "  - Measuring autonomous tool execution success rates and developer retention.",
    "What value is created for users and Redrob?",
    "  - Faster iteration cycles for Redrob, higher developer satisfaction."
])

# Slide 13: Future Roadmap
add_slide(prs, 1, "Future Roadmap", [
    "How could this evolve over 2-3 years?",
    "  - Transition from single-machine agents to distributed Swarm Intelligence managing entire cloud infrastructures.",
    "What future capabilities can be unlocked?",
    "  - Automated architectural migrations and proactive bug fixing before deployment.",
    "What broader vision does this support?",
    "  - Making Redrob the first fully AI-native company where systems build systems."
])

# Save the presentation
prs.save("Redrob_Hackathon_Pitch.pptx")
print("Successfully generated Redrob_Hackathon_Pitch.pptx")
