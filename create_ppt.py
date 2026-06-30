from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define some helper functions for styling
def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    return slide

def add_content_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        
    return slide

# Slide 1: Title Slide
add_title_slide(
    "CortexX AI System",
    "An AI-Native Autonomous Engineering System\nBuilt with Next.js & FastAPI"
)

# Slide 2: Overview
add_content_slide(
    "Project Overview",
    [
        "CortexX is a next-generation Multi-Agent AI system.",
        "Frontend: Stunning Next.js Glassmorphism UI.",
        "Backend: High-performance Python FastAPI server.",
        "LLMs: Powered by Google's Gemini advanced models.",
        "Two core modes: General Chat Copilot and Project Builder."
    ]
)

# Slide 3: General Chat Copilot
add_content_slide(
    "General Chat Copilot",
    [
        "Intelligent routing to specialized AI sub-agents.",
        "Code Intel Agent: Analyzes architecture and refactors code.",
        "Debug Agent: Rapid log parsing and root-cause identification.",
        "DevOps Agent: CI/CD pipeline monitoring and troubleshooting.",
        "Knowledge Agent: Semantic search across the codebase."
    ]
)

# Slide 4: Project Builder Mode
add_content_slide(
    "Autonomous Project Builder",
    [
        "Gives the AI physical agency over your local machine.",
        "Provides absolute workspace paths for code generation.",
        "LangChain ReAct loop executes tasks iteratively.",
        "Tool integrations: Read files, perform surgical inline edits.",
        "Executes terminal commands to install dependencies and run apps."
    ]
)

# Slide 5: Conclusion & Future
add_content_slide(
    "Conclusion",
    [
        "Bridges the gap between passive chat and active engineering.",
        "Acts as an autonomous digital teammate, not just a snippet generator.",
        "Secure local deployment with Cloudflare tunneling for remote access.",
        "Fully capable of building and maintaining complex systems from scratch.",
        "Revolutionizing how developers interact with their codebases."
    ]
)

# Save the presentation
prs.save("CortexX_Presentation.pptx")
print("Successfully generated CortexX_Presentation.pptx")
